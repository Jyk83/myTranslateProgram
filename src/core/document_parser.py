"""
문서 파싱 모듈
Excel, Word, PowerPoint, PDF 파일의 텍스트 추출 및 저장
"""

import os
from typing import Dict, List, Any, Optional
import traceback

# Excel 처리
try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
except ImportError:
    openpyxl = None

# Word 처리  
try:
    import docx
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    docx = None

# PowerPoint 처리
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    Presentation = None

# PDF 처리
try:
    import PyPDF2
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.lib.utils import ImageReader
    from io import BytesIO
except ImportError:
    PyPDF2 = None

from ..utils.logger import setup_logger


class DocumentContent:
    """문서 내용을 저장하는 클래스"""
    
    def __init__(self):
        self.content_type = ""  # excel, word, powerpoint, pdf
        self.original_format = ""
        self.text_content = []  # 텍스트 내용 리스트
        self.formatting_info = {}  # 포맷팅 정보
        self.metadata = {}  # 메타데이터
        
    def add_text_block(self, text, location_info=None, formatting=None):
        """텍스트 블록 추가"""
        self.text_content.append({
            'text': text,
            'location': location_info or {},
            'formatting': formatting or {}
        })


class DocumentParser:
    """문서 파싱 클래스"""
    
    def __init__(self):
        self.logger = setup_logger()
        self.supported_formats = ['.xlsx', '.docx', '.pptx', '.pdf']
        
    def parse_document(self, file_path: str) -> Optional[DocumentContent]:
        """문서 파싱 메인 함수"""
        try:
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")
                
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension not in self.supported_formats:
                raise ValueError(f"지원되지 않는 파일 형식: {file_extension}")
                
            # 파일 형식에 따른 파싱
            if file_extension == '.xlsx':
                return self.parse_excel(file_path)
            elif file_extension == '.docx':
                return self.parse_word(file_path)
            elif file_extension == '.pptx':
                return self.parse_powerpoint(file_path)
            elif file_extension == '.pdf':
                return self.parse_pdf(file_path)
            else:
                raise ValueError(f"지원되지 않는 파일 형식: {file_extension}")
                
        except Exception as e:
            self.logger.error(f"문서 파싱 오류: {e}")
            return None
            
    def parse_excel(self, file_path: str) -> Optional[DocumentContent]:
        """Excel 파일 파싱"""
        if not openpyxl:
            raise ImportError("openpyxl 라이브러리가 필요합니다.")
            
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            content = DocumentContent()
            content.content_type = "excel"
            content.original_format = ".xlsx"
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # 시트별로 텍스트 추출
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is not None and str(cell.value).strip():
                            # 수식이 아닌 텍스트만 추출
                            cell_text = str(cell.value).strip()
                            if cell_text and not cell_text.startswith('='):
                                location_info = {
                                    'sheet': sheet_name,
                                    'row': cell.row,
                                    'column': cell.column,
                                    'coordinate': cell.coordinate
                                }
                                
                                formatting = {
                                    'font_name': cell.font.name if cell.font else None,
                                    'font_size': cell.font.size if cell.font else None,
                                    'font_bold': cell.font.bold if cell.font else False,
                                    'font_italic': cell.font.italic if cell.font else False
                                }
                                
                                content.add_text_block(cell_text, location_info, formatting)
                                
            workbook.close()
            return content
            
        except Exception as e:
            self.logger.error(f"Excel 파싱 오류: {e}")
            return None
            
    def parse_word(self, file_path: str) -> Optional[DocumentContent]:
        """Word 문서 파싱"""
        if not docx:
            raise ImportError("python-docx 라이브러리가 필요합니다.")
            
        try:
            doc = Document(file_path)
            content = DocumentContent()
            content.content_type = "word"
            content.original_format = ".docx"
            
            # 문단별로 텍스트 추출
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    location_info = {
                        'paragraph_index': i,
                        'style': paragraph.style.name if paragraph.style else None
                    }
                    
                    formatting = {
                        'alignment': paragraph.alignment,
                        'runs': []
                    }
                    
                    # Run별 포맷팅 정보 수집
                    for run in paragraph.runs:
                        if run.text.strip():
                            run_formatting = {
                                'font_name': run.font.name,
                                'font_size': run.font.size.pt if run.font.size else None,
                                'bold': run.font.bold,
                                'italic': run.font.italic,
                                'underline': run.font.underline
                            }
                            formatting['runs'].append({
                                'text': run.text,
                                'formatting': run_formatting
                            })
                    
                    content.add_text_block(paragraph.text, location_info, formatting)
                    
            # 표 내용 추출
            for table_idx, table in enumerate(doc.tables):
                for row_idx, row in enumerate(table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            location_info = {
                                'table_index': table_idx,
                                'row_index': row_idx,
                                'cell_index': cell_idx,
                                'element_type': 'table_cell'
                            }
                            
                            content.add_text_block(cell.text.strip(), location_info)
                            
            return content
            
        except Exception as e:
            self.logger.error(f"Word 파싱 오류: {e}")
            return None
            
    def parse_powerpoint(self, file_path: str) -> Optional[DocumentContent]:
        """PowerPoint 파일 파싱"""
        if not Presentation:
            raise ImportError("python-pptx 라이브러리가 필요합니다.")
            
        try:
            prs = Presentation(file_path)
            content = DocumentContent()
            content.content_type = "powerpoint"
            content.original_format = ".pptx"
            
            for slide_idx, slide in enumerate(prs.slides):
                # 슬라이드의 모든 도형에서 텍스트 추출
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text") and shape.text.strip():
                        location_info = {
                            'slide_index': slide_idx,
                            'shape_index': shape_idx,
                            'shape_type': str(shape.shape_type)
                        }
                        
                        formatting = {
                            'shape_name': shape.name if hasattr(shape, 'name') else None
                        }
                        
                        # 텍스트 프레임이 있는 경우 상세 정보 추출
                        if hasattr(shape, "text_frame"):
                            text_frame = shape.text_frame
                            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                                if paragraph.text.strip():
                                    para_location = location_info.copy()
                                    para_location['paragraph_index'] = para_idx
                                    
                                    para_formatting = {
                                        'alignment': paragraph.alignment,
                                        'level': paragraph.level,
                                        'runs': []
                                    }
                                    
                                    for run in paragraph.runs:
                                        if run.text.strip():
                                            run_formatting = {
                                                'font_name': run.font.name,
                                                'font_size': run.font.size.pt if run.font.size else None,
                                                'bold': run.font.bold,
                                                'italic': run.font.italic
                                            }
                                            para_formatting['runs'].append({
                                                'text': run.text,
                                                'formatting': run_formatting
                                            })
                                    
                                    content.add_text_block(paragraph.text, para_location, para_formatting)
                        else:
                            content.add_text_block(shape.text, location_info, formatting)
                            
            return content
            
        except Exception as e:
            self.logger.error(f"PowerPoint 파싱 오류: {e}")
            return None
            
    def parse_pdf(self, file_path: str) -> Optional[DocumentContent]:
        """PDF 파일 파싱"""
        if not PyPDF2:
            raise ImportError("PyPDF2 라이브러리가 필요합니다.")
            
        try:
            content = DocumentContent()
            content.content_type = "pdf"
            content.original_format = ".pdf"
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_idx, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            location_info = {
                                'page_number': page_idx + 1,
                                'element_type': 'page_text'
                            }
                            
                            # 페이지별로 텍스트를 문단으로 분할
                            paragraphs = page_text.split('\n\n')
                            for para_idx, paragraph in enumerate(paragraphs):
                                if paragraph.strip():
                                    para_location = location_info.copy()
                                    para_location['paragraph_index'] = para_idx
                                    
                                    content.add_text_block(paragraph.strip(), para_location)
                                    
                    except Exception as e:
                        self.logger.warning(f"PDF 페이지 {page_idx + 1} 처리 오류: {e}")
                        continue
                        
            return content
            
        except Exception as e:
            self.logger.error(f"PDF 파싱 오류: {e}")
            return None
            
    def save_document(self, content: DocumentContent, output_path: str, original_format: str) -> Optional[str]:
        """번역된 문서를 원본 형식으로 저장"""
        try:
            if original_format == '.xlsx':
                return self.save_as_excel_original(content, output_path)
            elif original_format == '.docx':
                return self.save_as_word_original(content, output_path)
            elif original_format == '.pptx':
                return self.save_as_powerpoint_original(content, output_path)
            elif original_format == '.pdf':
                return self.save_as_pdf_original(content, output_path)
            else:
                raise ValueError(f"지원되지 않는 원본 형식: {original_format}")
                
        except Exception as e:
            self.logger.error(f"문서 저장 오류: {e}")
            return None
            
    def save_as_excel_original(self, content: DocumentContent, output_path: str) -> Optional[str]:
        """Excel 형식으로 저장 (원본 구조 유지)"""
        if not openpyxl:
            raise ImportError("openpyxl 라이브러리가 필요합니다.")
            
        try:
            workbook = openpyxl.Workbook()
            
            # 기본 워크시트 제거
            workbook.remove(workbook.active)
            
            # 시트별로 데이터 정리
            sheet_data = {}
            for text_block in content.text_content:
                location = text_block['location']
                sheet_name = location.get('sheet', 'Sheet1')
                
                if sheet_name not in sheet_data:
                    sheet_data[sheet_name] = []
                    
                sheet_data[sheet_name].append(text_block)
                
            # 시트별로 저장
            for sheet_name, blocks in sheet_data.items():
                worksheet = workbook.create_sheet(title=sheet_name)
                
                for block in blocks:
                    location = block['location']
                    row = location.get('row', 1)
                    column = location.get('column', 1)
                    
                    cell = worksheet.cell(row=row, column=column)
                    cell.value = block['text']
                    
                    # 포맷팅 적용
                    formatting = block.get('formatting', {})
                    if formatting.get('font_name') or formatting.get('font_size'):
                        cell.font = Font(
                            name=formatting.get('font_name', 'Calibri'),
                            size=formatting.get('font_size', 11),
                            bold=formatting.get('font_bold', False),
                            italic=formatting.get('font_italic', False)
                        )
                        
            workbook.save(output_path)
            return output_path
            
        except Exception as e:
            self.logger.error(f"Excel 저장 오류: {e}")
            return None
            
    def save_as_word_original(self, content: DocumentContent, output_path: str) -> Optional[str]:
        """Word 형식으로 저장 (원본 구조 유지)"""
        if not docx:
            raise ImportError("python-docx 라이브러리가 필요합니다.")
            
        try:
            doc = Document()
            
            # 문단별로 정리
            paragraphs = []
            tables = []
            
            for text_block in content.text_content:
                location = text_block['location']
                if location.get('element_type') == 'table_cell':
                    tables.append(text_block)
                else:
                    paragraphs.append(text_block)
                    
            # 문단 추가
            paragraphs.sort(key=lambda x: x['location'].get('paragraph_index', 0))
            for block in paragraphs:
                paragraph = doc.add_paragraph()
                
                # Run 정보가 있는 경우
                formatting = block.get('formatting', {})
                runs = formatting.get('runs', [])
                
                if runs:
                    for run_info in runs:
                        run = paragraph.add_run(run_info['text'])
                        run_fmt = run_info.get('formatting', {})
                        
                        if run_fmt.get('font_name'):
                            run.font.name = run_fmt['font_name']
                        if run_fmt.get('font_size'):
                            run.font.size = Pt(run_fmt['font_size'])
                        if run_fmt.get('bold'):
                            run.font.bold = True
                        if run_fmt.get('italic'):
                            run.font.italic = True
                else:
                    paragraph.add_run(block['text'])
                    
            # 테이블 처리 (간단한 형태로)
            if tables:
                table_data = {}
                for block in tables:
                    location = block['location']
                    table_idx = location.get('table_index', 0)
                    row_idx = location.get('row_index', 0)
                    cell_idx = location.get('cell_index', 0)
                    
                    if table_idx not in table_data:
                        table_data[table_idx] = {}
                    if row_idx not in table_data[table_idx]:
                        table_data[table_idx][row_idx] = {}
                        
                    table_data[table_idx][row_idx][cell_idx] = block['text']
                    
                for table_idx, table_rows in table_data.items():
                    max_cols = max(len(row) for row in table_rows.values()) if table_rows else 1
                    table = doc.add_table(rows=len(table_rows), cols=max_cols)
                    
                    for row_idx, row_data in table_rows.items():
                        for cell_idx, cell_text in row_data.items():
                            if row_idx < len(table.rows) and cell_idx < len(table.columns):
                                table.cell(row_idx, cell_idx).text = cell_text
                                
            doc.save(output_path)
            return output_path
            
        except Exception as e:
            self.logger.error(f"Word 저장 오류: {e}")
            return None
            
    def save_as_powerpoint_original(self, content: DocumentContent, output_path: str) -> Optional[str]:
        """PowerPoint 형식으로 저장 (원본 구조 유지)"""
        if not Presentation:
            raise ImportError("python-pptx 라이브러리가 필요합니다.")
            
        try:
            prs = Presentation()
            
            # 슬라이드별로 데이터 정리
            slide_data = {}
            for text_block in content.text_content:
                location = text_block['location']
                slide_idx = location.get('slide_index', 0)
                
                if slide_idx not in slide_data:
                    slide_data[slide_idx] = []
                    
                slide_data[slide_idx].append(text_block)
                
            # 슬라이드별로 처리
            for slide_idx in sorted(slide_data.keys()):
                slide_layout = prs.slide_layouts[1]  # 기본 레이아웃
                slide = prs.slides.add_slide(slide_layout)
                
                blocks = slide_data[slide_idx]
                
                # 제목과 내용 분리 (간단한 휴리스틱)
                title_blocks = []
                content_blocks = []
                
                for block in blocks:
                    location = block['location']
                    if location.get('shape_index', 0) == 0:  # 첫 번째 도형은 보통 제목
                        title_blocks.append(block)
                    else:
                        content_blocks.append(block)
                        
                # 제목 설정
                if title_blocks and slide.shapes.title:
                    slide.shapes.title.text = title_blocks[0]['text']
                    
                # 내용 설정
                if content_blocks and len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for block in content_blocks:
                        p = text_frame.add_paragraph()
                        p.text = block['text']
                        
            prs.save(output_path)
            return output_path
            
        except Exception as e:
            self.logger.error(f"PowerPoint 저장 오류: {e}")
            return None
            
    def save_as_pdf_original(self, content: DocumentContent, output_path: str) -> Optional[str]:
        """PDF 형식으로 저장"""
        return self.save_as_pdf(content, output_path)
        
    def save_as_pdf(self, content: DocumentContent, output_path: str) -> Optional[str]:
        """PDF로 변환 저장"""
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            
            # 한글 폰트 설정 시도
            try:
                # Windows의 경우 맑은 고딕 시도
                pdfmetrics.registerFont(TTFont('MalgunGothic', 'C:/Windows/Fonts/malgun.ttf'))
                font_name = 'MalgunGothic'
            except:
                try:
                    # 일반적인 한글 폰트 시도
                    pdfmetrics.registerFont(TTFont('NanumGothic', '/usr/share/fonts/truetype/nanum/NanumGothic.ttf'))
                    font_name = 'NanumGothic'
                except:
                    # 기본 폰트 사용
                    font_name = 'Helvetica'
                    
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            y_position = height - 50
            line_height = 20
            margin = 50
            
            for text_block in content.text_content:
                text = text_block['text']
                
                # 텍스트를 줄별로 분할하여 처리
                lines = text.split('\n')
                for line in lines:
                    if y_position < margin:
                        c.showPage()
                        y_position = height - 50
                        
                    c.setFont(font_name, 12)
                    
                    # 긴 줄을 페이지 너비에 맞게 분할
                    max_width = width - 2 * margin
                    if c.stringWidth(line, font_name, 12) > max_width:
                        words = line.split(' ')
                        current_line = ''
                        
                        for word in words:
                            test_line = current_line + ' ' + word if current_line else word
                            if c.stringWidth(test_line, font_name, 12) <= max_width:
                                current_line = test_line
                            else:
                                if current_line:
                                    c.drawString(margin, y_position, current_line)
                                    y_position -= line_height
                                    if y_position < margin:
                                        c.showPage()
                                        y_position = height - 50
                                current_line = word
                                
                        if current_line:
                            c.drawString(margin, y_position, current_line)
                            y_position -= line_height
                    else:
                        c.drawString(margin, y_position, line)
                        y_position -= line_height
                        
                y_position -= 10  # 블록 간 간격
                
            c.save()
            return output_path
            
        except Exception as e:
            self.logger.error(f"PDF 저장 오류: {e}")
            return None
            
    def save_as_excel(self, content: DocumentContent, output_path: str) -> Optional[str]:
        """Excel로 통합 저장"""
        if not openpyxl:
            raise ImportError("openpyxl 라이브러리가 필요합니다.")
            
        try:
            workbook = openpyxl.Workbook()
            worksheet = workbook.active
            worksheet.title = "번역 결과"
            
            # 헤더 설정
            worksheet['A1'] = "순번"
            worksheet['B1'] = "원본 위치"
            worksheet['C1'] = "번역된 텍스트"
            
            # 데이터 입력
            for i, text_block in enumerate(content.text_content, start=2):
                worksheet[f'A{i}'] = i - 1
                
                location = text_block['location']
                location_str = f"{location.get('sheet', location.get('slide_index', location.get('page_number', '')))} - {location.get('coordinate', location.get('paragraph_index', ''))}"
                worksheet[f'B{i}'] = location_str
                
                worksheet[f'C{i}'] = text_block['text']
                
            workbook.save(output_path)
            return output_path
            
        except Exception as e:
            self.logger.error(f"Excel 통합 저장 오류: {e}")
            return None
            
    def save_as_word(self, content: DocumentContent, output_path: str) -> Optional[str]:
        """Word로 통합 저장"""
        if not docx:
            raise ImportError("python-docx 라이브러리가 필요합니다.")
            
        try:
            doc = Document()
            
            # 제목 추가
            title = doc.add_heading('번역 결과', 0)
            
            # 내용 추가
            for i, text_block in enumerate(content.text_content, start=1):
                location = text_block['location']
                
                # 위치 정보 추가
                location_str = f"[{i}] 위치: {location.get('sheet', location.get('slide_index', location.get('page_number', 'N/A')))}"
                p = doc.add_paragraph()
                run = p.add_run(location_str)
                run.font.bold = True
                run.font.size = Pt(10)
                
                # 번역된 텍스트 추가
                text_para = doc.add_paragraph(text_block['text'])
                text_para.paragraph_format.left_indent = Inches(0.5)
                
                # 구분선
                doc.add_paragraph('─' * 50)
                
            doc.save(output_path)
            return output_path
            
        except Exception as e:
            self.logger.error(f"Word 통합 저장 오류: {e}")
            return None